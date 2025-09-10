"""Lógica de conversión PDF -> PPTX.

Responsabilidades:
- Exportar páginas de PDF a imágenes.
- Construir presentación PPTX con imágenes a pantalla completa (cover).
- Reportar progreso y logs mediante callbacks.
"""
from __future__ import annotations

import os
import time
from dataclasses import dataclass
from typing import Callable, List, Optional

try:
    import fitz  # PyMuPDF
except ImportError as e:  # pragma: no cover
    raise SystemExit("Falta 'pymupdf'. Instala con: pip install pymupdf") from e

from pptx import Presentation
from pptx.util import Inches

# --------------------------- Configuración --------------------------- #
IMG_DPI = 200

# --------------------------- Data Models --------------------------- #
@dataclass
class PdfConversionResult:
    pdf_path: str
    images: List[str]
    pptx_path: str
    first_page_px: tuple[int, int]
    dpi: int

# --------------------------- Utilidades internas --------------------------- #
def ensure_dir(path: str) -> None:
    os.makedirs(path, exist_ok=True)


def export_pdf_to_images(
    pdf_path: str,
    out_dir: str,
    dpi: int = IMG_DPI,
    log: Optional[Callable[[str], None]] = None,
    progress_hook: Optional[Callable[[float], None]] = None,
    progress_start: float = 0.02,
    progress_end: float = 0.50,
) -> List[str]:
    """Exporta cada página del PDF a JPG.

    Args:
        pdf_path: Ruta al PDF.
        out_dir: Carpeta destino de las imágenes.
        dpi: Resolución.
        log: Callback de log (opcional).
        progress_hook: Callback de progreso (opcional) 0..1.
        progress_start: Valor de progreso inicial reservado para esta fase.
        progress_end: Valor de progreso final reservado para esta fase.
    """
    ensure_dir(out_dir)
    doc = fitz.open(pdf_path)
    images: List[str] = []
    total_pages = len(doc)
    scale = dpi / 72
    matrix = fitz.Matrix(scale, scale)
    try:
        for page_index in range(total_pages):
            page = doc[page_index]
            pix = page.get_pixmap(matrix=matrix, alpha=False)
            img_name = f"page_{page_index + 1:03d}.jpg"
            img_path = os.path.join(out_dir, img_name)
            pix.save(img_path, jpg_quality=95)
            images.append(img_path)
            # Logging y progreso fino
            if log:
                log(f"Exportada página {page_index + 1}/{total_pages}")
            if progress_hook:
                frac = (page_index + 1) / total_pages
                progress_val = progress_start + (progress_end - progress_start) * frac
                progress_hook(progress_val)
    finally:
        doc.close()
    return images


def set_presentation_size_from_pixels(prs: Presentation, width_px: int, height_px: int, dpi: int) -> None:
    width_in = width_px / dpi
    height_in = height_px / dpi
    prs.slide_width = Inches(width_in)
    prs.slide_height = Inches(height_in)


def add_image_slide(prs: Presentation, image_path: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Layout en blanco
    picture = slide.shapes.add_picture(image_path, 0, 0)
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    img_w = picture.width
    img_h = picture.height
    scale = max(slide_w / img_w, slide_h / img_h)
    new_w = int(img_w * scale)
    new_h = int(img_h * scale)
    picture.width = new_w
    picture.height = new_h
    picture.left = int((slide_w - new_w) / 2)
    picture.top = int((slide_h - new_h) / 2)

# --------------------------- API Principal --------------------------- #
def convert_pdf_to_pptx(
    pdf_path: str,
    output_root: str,
    dpi: int = IMG_DPI,
    log: Optional[Callable[[str], None]] = None,
    progress: Optional[Callable[[float], None]] = None,
) -> PdfConversionResult:
    """Convierte un PDF a PPTX.

    Args:
        pdf_path: Ruta absoluta al PDF.
    output_root: Carpeta base (dentro se crea <nombre_pdf>/...)
        dpi: Resolución de rasterizado.
        log: Callback de logging.
        progress: Callback de progreso (0..1).
    """
    t0 = time.time()

    def _log(msg: str):
        if log:
            log(msg)

    def _progress(val: float):
        if progress:
            progress(max(0.0, min(1.0, val)))

    base = os.path.splitext(os.path.basename(pdf_path))[0]
    pdf_out_dir = os.path.join(output_root, base)
    img_dir = os.path.join(pdf_out_dir, "pages")

    _log(f"Abriendo PDF: {os.path.basename(pdf_path)}")
    _progress(0.02)
    images = export_pdf_to_images(
        pdf_path,
        img_dir,
        dpi=dpi,
        log=_log,
        progress_hook=_progress,
        progress_start=0.02,
        progress_end=0.50,  # ahora la exportación de páginas avanza hasta 50%
    )
    total_pages = len(images)
    if not images:
        raise RuntimeError(f"No se generaron imágenes para {pdf_path}")

    import PIL.Image as PILImage  # Lazy import
    with PILImage.open(images[0]) as im:
        first_w, first_h = im.width, im.height

    prs = Presentation()
    set_presentation_size_from_pixels(prs, first_w, first_h, dpi)

    _log("Añadiendo diapositivas...")
    for idx, img in enumerate(images, start=1):
        add_image_slide(prs, img)
        # Diapositivas: 50% -> 90%
        _progress(0.50 + 0.40 * (idx / total_pages))
        if idx % 1 == 0:  # siempre, pero fácil de ajustar si se quisiera espaciar
            _log(f"Diapositiva {idx}/{total_pages} creada")

    pptx_path = os.path.join(pdf_out_dir, base + ".pptx")
    ensure_dir(pdf_out_dir)
    _log("Guardando PPTX...")
    prs.save(pptx_path)
    _progress(0.95)
    dt = time.time() - t0
    _log(f"Terminado en {dt:.1f}s. Archivo: {pptx_path}")
    _progress(1.0)
    
    # Abrir la carpeta donde se guardó el resultado
    try:
        os.startfile(pdf_out_dir)
        _log(f"Carpeta abierta: {pdf_out_dir}")
    except Exception as e:
        _log(f"No se pudo abrir la carpeta: {e}")
    
    return PdfConversionResult(pdf_path, images, pptx_path, (first_w, first_h), dpi)
